#!/usr/bin/env python3
"""Generate a wafer yield comparison slide in an existing PPTX."""

import math
import random
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PPTX_PATH = "/Users/mike/Desktop/Mike/VSCode/Model-defined-Remapper/image/image.pptx"

# Slide dimensions in EMU
SLIDE_W = 12192000
SLIDE_H = 6858000

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
MED_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY = RGBColor(0x44, 0x44, 0x44)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREEN_FILL = RGBColor(0xC8, 0xE6, 0xC8)  # light green
GREEN_BORDER = RGBColor(0x5A, 0xA8, 0x5A)
RED_FILL = RGBColor(0xF2, 0xC4, 0xC4)  # light red
RED_BORDER = RGBColor(0xCC, 0x55, 0x55)
DEFECT_RED = RGBColor(0xCC, 0x22, 0x22)
WAFER_FILL = RGBColor(0xF0, 0xF0, 0xF0)
WAFER_BORDER = RGBColor(0x88, 0x88, 0x88)
BIG_RED_FILL = RGBColor(0xF2, 0xC4, 0xC4)


def emu(inches):
    return int(inches * 914400)


def add_textbox(slide, left, top, width, height, text, font_size=12,
                bold=False, color=DARK_GRAY, alignment=PP_ALIGN.CENTER):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def draw_circle(slide, cx, cy, radius, fill_color, border_color, border_width=Pt(1.5)):
    """Draw circle centered at (cx, cy) with given radius (all in EMU)."""
    left = cx - radius
    top = cy - radius
    diameter = 2 * radius
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, diameter, diameter
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    return shape


def draw_rect(slide, left, top, width, height, fill_color, border_color,
              border_width=Pt(0.75)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    return shape


def draw_defect_x(slide, cx, cy, size, color=DEFECT_RED):
    """Draw an X mark centered at (cx, cy)."""
    half = size // 2
    # Use a small filled circle as a defect marker (cleaner for academic look)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx - half, cy - half, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # no border
    return shape


def generate_defects(cx, cy, radius, n_defects=18, seed=42):
    """Generate random defect positions within a circle."""
    rng = random.Random(seed)
    defects = []
    while len(defects) < n_defects:
        # uniform in the disk
        r = radius * math.sqrt(rng.random()) * 0.88  # stay a bit inside
        theta = rng.random() * 2 * math.pi
        dx = int(r * math.cos(theta))
        dy = int(r * math.sin(theta))
        defects.append((cx + dx, cy + dy))
    return defects


def point_in_rect(px, py, left, top, width, height):
    return left <= px <= left + width and top <= py <= top + height


def point_in_circle(px, py, cx, cy, radius):
    return (px - cx) ** 2 + (py - cy) ** 2 <= radius ** 2


def main():
    prs = Presentation(PPTX_PATH)

    # Remove existing slide 2 if present (replace it)
    while len(prs.slides) > 1:
        rId = prs.slides._sldIdLst[-1].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])

    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # --- Layout constants ---
    margin_x = emu(0.6)
    panel_w = (SLIDE_W - 3 * margin_x) // 2
    panel_h = SLIDE_H - emu(2.0)  # leave room for title and bottom text

    left_panel_x = margin_x
    right_panel_x = 2 * margin_x + panel_w
    panel_top = emu(0.9)

    # Wafer circle parameters (relative to panel)
    wafer_radius = min(panel_w, panel_h) // 2 - emu(0.25)

    left_cx = left_panel_x + panel_w // 2
    left_cy = panel_top + panel_h // 2 - emu(0.1)

    right_cx = right_panel_x + panel_w // 2
    right_cy = panel_top + panel_h // 2 - emu(0.1)

    # --- Title ---
    add_textbox(slide, 0, emu(0.15), SLIDE_W, emu(0.55),
                "Wafer Yield vs. Die Size",
                font_size=24, bold=True, color=BLACK)

    # Generate defects in normalized coords [0,1] relative to wafer center
    defects_norm = []
    rng = random.Random(42)
    n_defects = 18
    while len(defects_norm) < n_defects:
        r = math.sqrt(rng.random()) * 0.85
        theta = rng.random() * 2 * math.pi
        defects_norm.append((r * math.cos(theta), r * math.sin(theta)))

    # Map to left/right panel absolute coords
    left_defects = [(left_cx + int(dx * wafer_radius),
                     left_cy + int(dy * wafer_radius))
                    for dx, dy in defects_norm]
    right_defects = [(right_cx + int(dx * wafer_radius),
                      right_cy + int(dy * wafer_radius))
                     for dx, dy in defects_norm]

    defect_marker_size = emu(0.08)
    label_y = left_cy + wafer_radius + emu(0.15)

    # Helper to draw a die grid panel
    # Only draw dies whose entire rectangle fits inside the wafer circle.
    def draw_die_panel(slide, cx, cy, grid_n, defects, panel_x, panel_w_):
        draw_circle(slide, cx, cy, wafer_radius, WAFER_FILL, WAFER_BORDER)
        die_gap = emu(0.03)
        # Fit the grid inside the inscribed square of the circle
        inscribed_half = int(wafer_radius / math.sqrt(2))  # half-side of inscribed square
        total_grid = 2 * inscribed_half - die_gap
        die_size = (total_grid - (grid_n - 1) * die_gap) // grid_n
        grid_origin_x = cx - total_grid // 2
        grid_origin_y = cy - total_grid // 2

        n_good = 0
        n_bad = 0
        for row in range(grid_n):
            for col in range(grid_n):
                dl = grid_origin_x + col * (die_size + die_gap)
                dt = grid_origin_y + row * (die_size + die_gap)
                # Check all 4 corners are inside the wafer circle
                corners = [
                    (dl, dt), (dl + die_size, dt),
                    (dl, dt + die_size), (dl + die_size, dt + die_size)
                ]
                if not all(point_in_circle(px, py, cx, cy, wafer_radius)
                           for px, py in corners):
                    continue
                hit = any(point_in_rect(dx, dy, dl, dt, die_size, die_size)
                          for dx, dy in defects)
                if hit:
                    draw_rect(slide, dl, dt, die_size, die_size,
                              RED_FILL, RED_BORDER)
                    n_bad += 1
                else:
                    draw_rect(slide, dl, dt, die_size, die_size,
                              GREEN_FILL, GREEN_BORDER)
                    n_good += 1

        for dx, dy in defects:
            draw_defect_x(slide, dx, dy, defect_marker_size)

        return n_good, n_bad

    # === LEFT PANEL: Large Dies (3x3 grid) ===
    lg_good, lg_bad = draw_die_panel(slide, left_cx, left_cy, 3,
                                      left_defects, left_panel_x, panel_w)
    lg_total = lg_good + lg_bad
    lg_yield = int(100 * lg_good / lg_total) if lg_total > 0 else 0
    add_textbox(slide, left_panel_x, label_y, panel_w, emu(0.45),
                f"Large Dies \u2014 {lg_yield}% Yield "
                f"({lg_good}/{lg_total} good)",
                font_size=16, bold=True, color=DARK_GRAY)

    # === RIGHT PANEL: Small Dies (10x10 grid) ===
    sm_good, sm_bad = draw_die_panel(slide, right_cx, right_cy, 10,
                                      right_defects, right_panel_x, panel_w)
    sm_total = sm_good + sm_bad
    sm_yield = int(100 * sm_good / sm_total) if sm_total > 0 else 0
    add_textbox(slide, right_panel_x, label_y, panel_w, emu(0.45),
                f"Small Dies \u2014 {sm_yield}% Yield "
                f"({sm_good}/{sm_total} good)",
                font_size=16, bold=True, color=DARK_GRAY)

    # === Bottom text ===
    bottom_y = SLIDE_H - emu(0.65)
    add_textbox(slide, emu(0.5), bottom_y, SLIDE_W - emu(1.0), emu(0.5),
                "Same defect distribution \u2192 larger dies have lower yield; "
                "at wafer scale, defect-free fabrication is impossible.",
                font_size=14, bold=False, color=MED_GRAY)

    # Save
    prs.save(PPTX_PATH)
    print(f"Saved to {PPTX_PATH}")
    print(f"Large dies: {lg_good}/{lg_total} good ({lg_yield}% yield)")
    print(f"Small dies: {sm_good}/{sm_total} good ({sm_yield}% yield)")


if __name__ == "__main__":
    main()
